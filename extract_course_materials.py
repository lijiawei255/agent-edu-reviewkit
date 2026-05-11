"""
课程课件内容提取工具 —— 文本 + 图片
支持格式: .pdf (pypdf), .pptx/.pptm (python-pptx), .docx/.dotx/.dotm (python-docx + zipfile)
所有依赖均为纯Python库，pip一键安装，无C编译依赖。

使用方法:
  1. 修改下方"配置区"的 课件目录、文本输出目录、图片输出目录
  2. 运行: python extract_course_materials.py
"""

import os
import sys
import zipfile
from pathlib import Path

# 修复Windows控制台编码问题（GBK无法输出Unicode字符如 ✓ ✗）
if sys.platform == 'win32':
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

# ============================================================
# 0. 依赖自检
# ============================================================
MISSING = []
try:
    from pypdf import PdfReader
except ImportError:
    MISSING.append("pypdf")
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    MISSING.append("python-pptx")
try:
    from docx import Document
except ImportError:
    MISSING.append("python-docx")

if MISSING:
    print("=" * 50)
    print("以下 Python 库未安装，请先运行：")
    print(f"  pip install {' '.join(MISSING)}")
    print("=" * 50)
    sys.exit(1)

# ============================================================
# 1. 配置区 —— 用户根据需要修改以下变量
# ============================================================
课件目录 = r"课件"               # 课件所在目录路径
文本输出目录 = r"extracted_text"   # 提取文本输出目录
图片输出目录 = r"extracted_images" # 提取图片输出目录

os.makedirs(文本输出目录, exist_ok=True)
os.makedirs(图片输出目录, exist_ok=True)

# ============================================================
# 2. 工具函数
# ============================================================

def _content_type_to_ext(content_type):
    """将 MIME content-type 映射为文件扩展名"""
    mapping = {
        "image/png":  ".png",
        "image/jpeg": ".jpg",
        "image/gif":  ".gif",
        "image/bmp":  ".bmp",
        "image/tiff": ".tiff",
        "image/webp": ".webp",
        "image/x-png": ".png",
        "image/x-emf": ".emf",
        "image/x-wmf": ".wmf",
    }
    return mapping.get(content_type, ".png")


def _guess_ext_from_bytes(data):
    """通过文件头魔术字节推断图片扩展名"""
    if data[:4] == b'\x89PNG':
        return '.png'
    if data[:2] == b'\xff\xd8':
        return '.jpg'
    if data[:4] == b'GIF8':
        return '.gif'
    if data[:2] == b'BM':
        return '.bmp'
    if data[:4] == b'RIFF' and data[8:12] == b'WEBP':
        return '.webp'
    return '.png'


def _sanitize_filename(s):
    """移除文件名中的非法字符"""
    illegal = '<>:"/\\|?*'
    for c in illegal:
        s = s.replace(c, '_')
    return s


# ============================================================
# 3. 纯图片文件预检
# ============================================================

def _check_likely_image_only(filepath, ext):
    """预检文件是否可能是纯图片型（扫描版PDF等）。
    检查前三页的文本量：如果不足50个字符，大概率是扫描版/图片型PDF。
    对于PPTX/DOCX，不做预检（它们几乎总有文字，即使很少）。"""
    if ext == '.pdf':
        try:
            reader = PdfReader(filepath)
            if len(reader.pages) == 0:
                return False
            total_chars = 0
            pages_to_check = min(3, len(reader.pages))
            for page in reader.pages[:pages_to_check]:
                text = page.extract_text()
                if text:
                    total_chars += len(text.strip())
            return total_chars < 50
        except Exception:
            return False
    return False


# ============================================================
# 4. 格式提取函数 (文本 + 图片)
# ============================================================

def extract_pdf(filepath, basename):
    """提取 PDF 文本和图片"""
    reader = PdfReader(filepath)
    lines = []
    img_count = 0

    for i, page in enumerate(reader.pages):
        page_num = i + 1

        # --- 文本 ---
        text = page.extract_text()
        if text and text.strip():
            lines.append(f"\n--- 第{page_num}页 ---")
            lines.append(text)

        # --- 图片 ---
        try:
            for j, img_obj in enumerate(page.images):
                img_data = img_obj.data
                if not img_data or len(img_data) < 100:
                    continue
                ext = _guess_ext_from_bytes(img_data)
                img_name = f"{basename}_p{page_num}_img{j+1}{ext}"
                img_path = os.path.join(图片输出目录, img_name)
                with open(img_path, "wb") as f:
                    f.write(img_data)
                lines.append(f"[图片: {img_name}]")
                img_count += 1
        except Exception:
            pass

    return "\n".join(lines), len(reader.pages), img_count


def extract_pptx(filepath, basename):
    """提取 PPTX/PPTM 文本和图片"""
    prs = Presentation(filepath)
    lines = []
    img_count = 0

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_texts = []

        for shape in slide.shapes:
            # --- 图片 ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    if img_blob and len(img_blob) >= 100:
                        ext = _content_type_to_ext(shape.image.content_type)
                        img_name = f"{basename}_s{slide_num}_img{img_count+1}{ext}"
                        img_path = os.path.join(图片输出目录, img_name)
                        with open(img_path, "wb") as f:
                            f.write(img_blob)
                        slide_texts.append(f"[图片: {img_name}]")
                        img_count += 1
                except Exception:
                    pass

            # --- 文本框 ---
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)

            # --- 表格 ---
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        ct = cell.text.strip().replace("\n", " ")
                        row_texts.append(ct)
                    slide_texts.append(" | ".join(row_texts))

        if slide_texts:
            lines.append(f"\n--- 第{slide_num}张幻灯片 ---")
            lines.append("\n".join(slide_texts))

    return "\n".join(lines), len(prs.slides), img_count


def extract_docx(filepath, basename):
    """提取 DOCX/DOTX/DOTM 文本和图片 (图片通过 zipfile 从 word/media/ 提取)"""
    doc = Document(filepath)
    lines = []
    img_count = 0

    # --- 文本 (通过 python-docx) ---
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            lines.append(t)

    for ti, table in enumerate(doc.tables):
        lines.append(f"\n[表格 {ti+1}]")
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                ct = cell.text.strip().replace("\n", " ")
                row_texts.append(ct)
            lines.append(" | ".join(row_texts))

    # --- 图片 (通过 zipfile 从 docx 内部提取) ---
    try:
        with zipfile.ZipFile(filepath) as z:
            media_files = [n for n in z.namelist() if n.startswith("word/media/")]
            for idx, media_path in enumerate(media_files):
                img_data = z.read(media_path)
                if not img_data or len(img_data) < 100:
                    continue
                original_name = os.path.basename(media_path)
                ext = os.path.splitext(original_name)[1].lower()
                if ext not in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'):
                    ext = _guess_ext_from_bytes(img_data)
                img_name = f"{basename}_img{idx+1}{ext}"
                img_path = os.path.join(图片输出目录, img_name)
                with open(img_path, "wb") as f:
                    f.write(img_data)
                lines.append(f"[图片: {img_name}]")
                img_count += 1
    except Exception:
        pass

    return "\n".join(lines), len(doc.paragraphs), img_count


# ============================================================
# 5. 扫描并提取
# ============================================================
HANDLERS = {
    ".pdf":  extract_pdf,
    ".pptx": extract_pptx,
    ".pptm": extract_pptx,
    ".docx": extract_docx,
    ".dotx": extract_docx,
    ".dotm": extract_docx,
}

目标文件 = []
for fname in sorted(os.listdir(课件目录)):
    ext = os.path.splitext(fname)[1].lower()
    if ext in HANDLERS:
        目标文件.append(fname)

if not 目标文件:
    print(f"在 '{课件目录}' 中未找到支持的课件文件，请检查路径。")
    print(f"支持格式: {', '.join(sorted(set(HANDLERS.keys())))}")
    sys.exit(1)

print(f"找到 {len(目标文件)} 个课件文件，开始提取文本和图片...\n")

总图片数 = 0
成功 = 0
失败 = []
空内容 = []
纯图片文件 = []

for idx, fname in enumerate(目标文件, 1):
    filepath = os.path.join(课件目录, fname)
    ext = os.path.splitext(fname)[1].lower()
    basename = _sanitize_filename(Path(fname).stem)
    handler = HANDLERS[ext]

    # 预检：是否为纯图片文件
    if _check_likely_image_only(filepath, ext):
        纯图片文件.append(fname)
        print(f"  🔍 [{idx}/{len(目标文件)}] {fname} — 检测为纯图片型文件（扫描版PDF等），跳过文本提取")
        continue

    try:
        text, page_count, img_count = handler(filepath, basename)
        header = f"======= {fname} =======\n页数/幻灯片数: {page_count}\n" + "=" * 60

        if text.strip():
            outpath = os.path.join(文本输出目录, fname + ".txt")
            with open(outpath, "w", encoding="utf-8") as f:
                f.write(header + "\n" + text)
            成功 += 1
            总图片数 += img_count
            img_info = f", {img_count}张图片" if img_count else ""
            print(f"  ✓ [{成功}/{len(目标文件)}] {fname} ({page_count}p{img_info})")
        else:
            空内容.append(fname)
            print(f"  ⚠ [{成功}/{len(目标文件)}] {fname} — 提取内容为空")
    except Exception as e:
        失败.append((fname, str(e)))
        print(f"  ✗ [{成功}/{len(目标文件)}] {fname} — 错误: {e}")

# ============================================================
# 6. 结果报告
# ============================================================
print(f"\n{'=' * 50}")
print(f"提取完成: 成功 {成功}/{len(目标文件)}, 共提取 {总图片数} 张图片")
print(f"  文本目录: {os.path.abspath(文本输出目录)}")
print(f"  图片目录: {os.path.abspath(图片输出目录)}")

if 纯图片文件:
    print(f"\n🔍 以下 {len(纯图片文件)} 个文件被检测为纯图片型（扫描版PDF等）：")
    for fn in 纯图片文件:
        print(f"  - {fn}")
    print("这些文件需要通过AI视觉能力直接读取，不依赖文本提取。")
    print("如果你的AI助手不支持图片识别，请参考以下方案：")
    print("  方案A: 更换为支持视觉的模型（如 Claude Opus 4、GPT-4V、Gemini Pro Vision）")
    print("  方案B: 安装MCP视觉服务器（如 MiniMax MCP）")
    print("  方案C: 使用本地OCR工具预处理：pip install pytesseract pdf2image")

if 空内容:
    print(f"\n⚠ 以下文件提取内容为空 ({len(空内容)}个):")
    for fn in 空内容:
        print(f"  - {fn}")
    print("可能原因：PPTX/DOCX只含图片不含文字。")
    print("建议：请AI助手通过视觉能力直接读取原始文件中的图片内容。")

if 失败:
    print(f"\n✗ 以下文件提取失败 ({len(失败)}个):")
    for fn, err in 失败:
        print(f"  - {fn}: {err}")
    print("建议：检查原始文件是否损坏，或尝试手动转换格式后重新提取。")

if not 空内容 and not 失败 and not 纯图片文件:
    print("所有文件提取成功！")
